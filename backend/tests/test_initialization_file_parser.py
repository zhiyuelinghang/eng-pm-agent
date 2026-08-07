import io
import importlib.util
from pathlib import Path

import pytest
from docx import Document
from openpyxl import Workbook
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches


PARSER_PATH = (
    Path(__file__).resolve().parents[2]
    / "mcp-packages"
    / "attachment-parser"
    / "local_parser.py"
)
PARSER_SPEC = importlib.util.spec_from_file_location(
    "initialization_attachment_local_parser",
    PARSER_PATH,
)
assert PARSER_SPEC is not None and PARSER_SPEC.loader is not None
initialization_file_parser = importlib.util.module_from_spec(PARSER_SPEC)
PARSER_SPEC.loader.exec_module(initialization_file_parser)
_parse_attachment_content = initialization_file_parser.parse_attachment_content


SAMPLES = (
    (
        "原型/工程基本信息/1、工程描述.txt",
        "txt",
        "total_lines",
    ),
    (
        "原型/工程基本信息/2、总进度计划（WBS）.xlsx",
        "xlsx",
        "total_rows",
    ),
    (
        "原型/工程基本信息/3.人员名单 - 副本.xlsx",
        "xlsx",
        "total_rows",
    ),
    (
        "原型/工程基本信息/4、工程风险清单.xlsx",
        "xlsx",
        "total_rows",
    ),
    (
        "原型/工程基本信息/5.工序质量指标关联表.xlsx",
        "xlsx",
        "total_rows",
    ),
)


@pytest.mark.parametrize(("relative_path", "expected_format", "count_key"), SAMPLES)
def test_project_initialization_samples_can_be_parsed(
    relative_path: str,
    expected_format: str,
    count_key: str,
) -> None:
    path = Path(relative_path)
    result = _parse_attachment_content(
        path.read_bytes(),
        path.suffix.lower(),
        None,
        1,
        10,
    )

    assert result["format"] == expected_format
    assert result[count_key] > 0


def test_wbs_excel_keeps_zero_progress_distinct_from_empty_cells() -> None:
    path = Path(
        "原型/工程基本信息/2、总进度计划（WBS）.xlsx",
    )
    result = _parse_attachment_content(
        path.read_bytes(),
        path.suffix.lower(),
        None,
        5,
        3,
    )

    header, root_wbs, child_wbs = result["rows"]
    progress_column = header.index("进度 (%)")
    assert root_wbs[progress_column] == 0
    assert child_wbs[progress_column] == 0


def test_xlsx_returns_merge_formula_and_zero_metadata() -> None:
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "工程"
    worksheet["A1"] = "项目"
    worksheet.merge_cells("A1:B1")
    worksheet["A2"] = 0
    worksheet["B2"] = None
    worksheet["C2"] = "=A2+1"
    buffer = io.BytesIO()
    workbook.save(buffer)

    result = _parse_attachment_content(
        buffer.getvalue(),
        ".xlsx",
        "工程",
        1,
        10,
        "never",
    )

    assert result["rows"][1][0] == 0
    assert result["rows"][1][1] is None
    assert result["rows"][1][2] == "=A2+1"
    assert result["merged_ranges"] == ["A1:B1"]
    assert result["formulas"] == [
        {
            "cell": "C2",
            "formula": "=A2+1",
            "cached_value": None,
        },
    ]


def test_docx_blocks_keep_paragraph_and_table_order() -> None:
    document = Document()
    document.add_paragraph("第一段")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "字段"
    table.cell(0, 1).text = "值"
    document.add_paragraph("第二段")
    buffer = io.BytesIO()
    document.save(buffer)

    result = _parse_attachment_content(
        buffer.getvalue(),
        ".docx",
        None,
        1,
        10,
        "never",
    )

    assert [block["type"] for block in result["blocks"]] == [
        "paragraph",
        "table",
        "paragraph",
    ]
    assert result["blocks"][1]["rows"] == [["字段", "值"]]


def test_pptx_returns_slide_text_and_tables() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(
        presentation.slide_layouts[5],
    )
    slide.shapes.title.text = "项目汇报"
    table = slide.shapes.add_table(
        2,
        2,
        Inches(1),
        Inches(2),
        Inches(5),
        Inches(2),
    ).table
    table.cell(0, 0).text = "指标"
    table.cell(0, 1).text = "结果"
    table.cell(1, 0).text = "进度"
    table.cell(1, 1).text = "0"
    buffer = io.BytesIO()
    presentation.save(buffer)

    result = _parse_attachment_content(
        buffer.getvalue(),
        ".pptx",
        None,
        1,
        10,
        "never",
    )

    assert result["total_slides"] == 1
    assert "项目汇报" in result["slides"][0]["texts"]
    assert result["slides"][0]["tables"][0]["rows"][1] == [
        "进度",
        "0",
    ]


def test_scanned_pdf_uses_ocr_in_auto_mode(monkeypatch) -> None:
    image = Image.new("RGB", (300, 160), "white")
    buffer = io.BytesIO()
    image.save(buffer, format="PDF")
    calls: list[str] = []

    monkeypatch.setattr(
        initialization_file_parser,
        "_render_pdf_page",
        lambda content, page_index: b"rendered",
    )

    def fake_ocr(content: bytes, *, source: str) -> dict[str, object]:
        calls.append(source)
        return {"source": source, "text": "扫描工程资料", "lines": []}

    monkeypatch.setattr(
        initialization_file_parser,
        "_ocr_image",
        fake_ocr,
    )

    result = _parse_attachment_content(
        buffer.getvalue(),
        ".pdf",
        None,
        1,
        10,
        "auto",
    )

    assert calls == ["PDF 第1页"]
    assert result["pages"][0]["ocr"]["text"] == "扫描工程资料"


def test_image_attachment_can_be_ocr_recognized() -> None:
    image = Image.new("RGB", (1000, 260), "white")
    draw = ImageDraw.Draw(image)
    font = ImageFont.load_default(size=80)
    draw.text((40, 70), "DOBBY 123", fill="black", font=font)
    buffer = io.BytesIO()
    image.save(buffer, format="PNG")

    result = _parse_attachment_content(
        buffer.getvalue(),
        ".png",
        None,
        1,
        10,
        "auto",
    )

    recognized = result["image"]["ocr"]["text"].upper()
    assert "DOBBY" in recognized
    assert "123" in recognized
