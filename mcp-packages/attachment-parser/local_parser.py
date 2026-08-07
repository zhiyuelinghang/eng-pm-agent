"""Deterministic fallback parsers bundled inside the attachment parser MCP.

The MCP tries the configured MinerU service first for the open-source input
formats.  These bounded local parsers preserve the existing offline behavior
when MinerU is unavailable or rejects a document.
"""

from __future__ import annotations

import csv
import io
import threading
from collections.abc import Iterable
from datetime import date, datetime
from typing import Any, Literal


OCRMode = Literal["auto", "always", "never"]

SUPPORTED_ATTACHMENT_SUFFIXES = frozenset(
    {
        ".txt",
        ".md",
        ".csv",
        ".xls",
        ".xlsx",
        ".docx",
        ".pptx",
        ".pdf",
        ".png",
        ".jpg",
        ".jpeg",
        ".bmp",
        ".webp",
        ".gif",
        ".jp2",
        ".tif",
        ".tiff",
    },
)

_IMAGE_SUFFIXES = frozenset(
    {
        ".png",
        ".jpg",
        ".jpeg",
        ".bmp",
        ".webp",
        ".gif",
        ".jp2",
        ".tif",
        ".tiff",
    },
)
_MAX_OCR_EDGE = 3200
_MAX_OCR_IMAGES_PER_CALL = 20
_MAX_TABLES_PER_PAGE = 20
_MAX_TABLE_ROWS = 500
_MAX_CELL_TEXT = 20_000
_OCR_ENGINE: Any = None
_OCR_ENGINE_LOCK = threading.Lock()


def _json_cell(value: Any) -> Any:
    """Return one JSON-safe cell while preserving numeric zero."""
    if value is None or isinstance(value, (str, int, float, bool)):
        return value
    if isinstance(value, (datetime, date)):
        return value.isoformat()
    if hasattr(value, "isoformat"):
        return value.isoformat()
    return str(value)


def _decode_text_file(content: bytes) -> tuple[str, str]:
    """Decode common Chinese project text files and return the encoding."""
    for encoding in ("utf-8-sig", "utf-8", "gb18030"):
        try:
            return content.decode(encoding), encoding
        except UnicodeDecodeError:
            continue
    return content.decode("utf-8", errors="replace"), "utf-8-replace"


def _normalize_ocr_mode(value: str | None) -> OCRMode:
    normalized = (value or "auto").strip().lower()
    if normalized not in {"auto", "always", "never"}:
        raise ValueError("ocr_mode 仅支持 auto、always 或 never。")
    return normalized  # type: ignore[return-value]


def _get_ocr_engine() -> Any:
    """Create the local RapidOCR engine lazily and reuse it safely."""
    global _OCR_ENGINE
    with _OCR_ENGINE_LOCK:
        if _OCR_ENGINE is not None:
            return _OCR_ENGINE
        try:
            from rapidocr_onnxruntime import RapidOCR
        except ImportError as exc:
            raise ValueError(
                "本地 OCR 组件未安装，请在项目内嵌 Python 中安装 "
                "rapidocr-onnxruntime。",
            ) from exc
        _OCR_ENGINE = RapidOCR()
        return _OCR_ENGINE


def _ocr_image(
    content: bytes,
    *,
    source: str,
) -> dict[str, Any]:
    """OCR one image and return bounded text, confidence and coordinates."""
    import numpy as np
    from PIL import Image, ImageOps

    try:
        with Image.open(io.BytesIO(content)) as original:
            image = ImageOps.exif_transpose(original).convert("RGB")
            original_width, original_height = image.size
            if max(image.size) > _MAX_OCR_EDGE:
                image.thumbnail(
                    (_MAX_OCR_EDGE, _MAX_OCR_EDGE),
                    Image.Resampling.LANCZOS,
                )
            width, height = image.size
            image_array = np.asarray(image)
    except Exception as exc:
        raise ValueError(f"无法打开待识别图片 {source}：{exc}") from exc

    engine = _get_ocr_engine()
    with _OCR_ENGINE_LOCK:
        result, _ = engine(image_array)
    lines: list[dict[str, Any]] = []
    for item in result or []:
        if not isinstance(item, (list, tuple)) or len(item) < 3:
            continue
        box, text, confidence = item[0], str(item[1]), item[2]
        points = [
            [round(float(point[0]), 2), round(float(point[1]), 2)]
            for point in box
        ]
        lines.append(
            {
                "text": text,
                "confidence": round(float(confidence), 4),
                "box": points,
            },
        )
    return {
        "source": source,
        "width": width,
        "height": height,
        "original_width": original_width,
        "original_height": original_height,
        "text": "\n".join(line["text"] for line in lines),
        "lines": lines,
    }


def _bounded_rows(
    rows: Iterable[Iterable[Any]],
    *,
    limit: int = _MAX_TABLE_ROWS,
) -> list[list[Any]]:
    result: list[list[Any]] = []
    for row in rows:
        if len(result) >= limit:
            break
        result.append(
            [
                (
                    cell[:_MAX_CELL_TEXT]
                    if isinstance(cell, str)
                    else _json_cell(cell)
                )
                for cell in row
            ],
        )
    return result


def _parse_xlsx(
    content: bytes,
    sheet_name: str | None,
    start: int,
    limit: int,
    ocr_mode: OCRMode,
) -> dict[str, Any]:
    from openpyxl import load_workbook

    values_workbook = load_workbook(
        io.BytesIO(content),
        read_only=False,
        data_only=True,
    )
    formulas_workbook = load_workbook(
        io.BytesIO(content),
        read_only=False,
        data_only=False,
    )
    try:
        sheets = list(values_workbook.sheetnames)
        selected_sheet = sheet_name or (sheets[0] if sheets else None)
        if selected_sheet is None or selected_sheet not in sheets:
            raise ValueError(
                f"工作表不存在，可选工作表：{sheets}",
            )
        values_sheet = values_workbook[selected_sheet]
        formulas_sheet = formulas_workbook[selected_sheet]
        total_rows = values_sheet.max_row
        end = min(total_rows, start + limit - 1)

        rows: list[list[Any]] = []
        formulas: list[dict[str, Any]] = []
        number_formats: list[dict[str, Any]] = []
        for row_number in range(start, end + 1):
            row_values: list[Any] = []
            for column_number in range(1, values_sheet.max_column + 1):
                value_cell = values_sheet.cell(row_number, column_number)
                formula_cell = formulas_sheet.cell(row_number, column_number)
                value = value_cell.value
                formula = formula_cell.value
                if (
                    isinstance(formula, str)
                    and formula.startswith("=")
                ):
                    formulas.append(
                        {
                            "cell": formula_cell.coordinate,
                            "formula": formula,
                            "cached_value": _json_cell(value),
                        },
                    )
                    if value is None:
                        value = formula
                row_values.append(_json_cell(value))
                if formula_cell.number_format != "General":
                    number_formats.append(
                        {
                            "cell": formula_cell.coordinate,
                            "format": formula_cell.number_format,
                        },
                    )
            rows.append(row_values)

        merged_ranges = []
        for merged_range in formulas_sheet.merged_cells.ranges:
            if (
                merged_range.max_row >= start
                and merged_range.min_row <= end
            ):
                merged_ranges.append(str(merged_range))

        images: list[dict[str, Any]] = []
        if ocr_mode != "never":
            for image_index, image in enumerate(
                getattr(formulas_sheet, "_images", []),
                start=1,
            ):
                if len(images) >= _MAX_OCR_IMAGES_PER_CALL:
                    break
                try:
                    anchor = image.anchor._from
                    anchor_cell = (
                        f"{formulas_sheet.cell(anchor.row + 1, anchor.col + 1).coordinate}"
                    )
                    image_result = _ocr_image(
                        image._data(),
                        source=f"{selected_sheet}!{anchor_cell} 图片{image_index}",
                    )
                    image_result["anchor_cell"] = anchor_cell
                    images.append(image_result)
                except Exception as exc:
                    images.append(
                        {
                            "source": f"{selected_sheet} 图片{image_index}",
                            "error": str(exc),
                        },
                    )

        return {
            "format": "xlsx",
            "sheet_names": sheets,
            "sheet_name": selected_sheet,
            "start_row": start,
            "end_row": end,
            "total_rows": total_rows,
            "total_columns": values_sheet.max_column,
            "rows": rows,
            "merged_ranges": merged_ranges,
            "formulas": formulas,
            "number_formats": number_formats,
            "images": images,
            "next_start": end + 1 if end < total_rows else None,
        }
    finally:
        values_workbook.close()
        formulas_workbook.close()


def _parse_xls(
    content: bytes,
    sheet_name: str | None,
    start: int,
    limit: int,
) -> dict[str, Any]:
    import xlrd

    workbook = xlrd.open_workbook(file_contents=content)
    sheets = workbook.sheet_names()
    selected_sheet = sheet_name or (sheets[0] if sheets else None)
    if selected_sheet is None or selected_sheet not in sheets:
        raise ValueError(f"工作表不存在，可选工作表：{sheets}")
    worksheet = workbook.sheet_by_name(selected_sheet)
    end = min(worksheet.nrows, start + limit - 1)
    rows: list[list[Any]] = []
    for row_index in range(start - 1, end):
        row: list[Any] = []
        for cell in worksheet.row(row_index):
            value = cell.value
            if cell.ctype == xlrd.XL_CELL_DATE:
                value = xlrd.xldate_as_datetime(value, workbook.datemode)
            elif cell.ctype == xlrd.XL_CELL_BOOLEAN:
                value = bool(value)
            elif cell.ctype == xlrd.XL_CELL_EMPTY:
                value = None
            row.append(_json_cell(value))
        rows.append(row)
    merged_ranges = [
        {
            "start_row": row_start + 1,
            "end_row": row_end,
            "start_column": column_start + 1,
            "end_column": column_end,
        }
        for row_start, row_end, column_start, column_end
        in worksheet.merged_cells
        if row_end >= start and row_start + 1 <= end
    ]
    return {
        "format": "xls",
        "sheet_names": sheets,
        "sheet_name": selected_sheet,
        "start_row": start,
        "end_row": end,
        "total_rows": worksheet.nrows,
        "total_columns": worksheet.ncols,
        "rows": rows,
        "merged_ranges": merged_ranges,
        "next_start": end + 1 if end < worksheet.nrows else None,
    }


def _docx_picture_blocks(
    paragraph_element: Any,
    document: Any,
    ocr_mode: OCRMode,
    image_counter: list[int],
) -> list[dict[str, Any]]:
    if ocr_mode == "never":
        return []
    from docx.oxml.ns import qn

    blocks: list[dict[str, Any]] = []
    for blip in paragraph_element.xpath(".//a:blip"):
        if image_counter[0] >= _MAX_OCR_IMAGES_PER_CALL:
            break
        relation_id = blip.get(qn("r:embed"))
        if not relation_id:
            continue
        image_counter[0] += 1
        try:
            part = document.part.related_parts[relation_id]
            ocr = _ocr_image(
                part.blob,
                source=f"DOCX 图片{image_counter[0]}",
            )
            blocks.append({"type": "image", "ocr": ocr})
        except Exception as exc:
            blocks.append(
                {
                    "type": "image",
                    "source": f"DOCX 图片{image_counter[0]}",
                    "error": str(exc),
                },
            )
    return blocks


def _parse_docx(
    content: bytes,
    start: int,
    limit: int,
    ocr_mode: OCRMode,
) -> dict[str, Any]:
    from docx import Document
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    document = Document(io.BytesIO(content))
    blocks: list[dict[str, Any]] = []
    table_index = 0
    image_counter = [0]
    for child in document.element.body.iterchildren():
        if isinstance(child, CT_P):
            paragraph = Paragraph(child, document)
            if paragraph.text.strip():
                blocks.append(
                    {
                        "type": "paragraph",
                        "text": paragraph.text[:_MAX_CELL_TEXT],
                    },
                )
            blocks.extend(
                _docx_picture_blocks(
                    child,
                    document,
                    ocr_mode,
                    image_counter,
                ),
            )
        elif isinstance(child, CT_Tbl):
            table_index += 1
            table = Table(child, document)
            blocks.append(
                {
                    "type": "table",
                    "table": table_index,
                    "rows": _bounded_rows(
                        (
                            (cell.text for cell in row.cells)
                            for row in table.rows
                        ),
                    ),
                },
            )

    selected = blocks[start - 1:start - 1 + limit]
    end = start + len(selected) - 1
    return {
        "format": "docx",
        "start_block": start,
        "end_block": end,
        "total_blocks": len(blocks),
        "blocks": selected,
        "next_start": end + 1 if end < len(blocks) else None,
    }


def _iter_presentation_shapes(shapes: Iterable[Any]) -> Iterable[Any]:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in sorted(
        shapes,
        key=lambda item: (
            int(getattr(item, "top", 0)),
            int(getattr(item, "left", 0)),
        ),
    ):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_presentation_shapes(shape.shapes)
        else:
            yield shape


def _parse_pptx(
    content: bytes,
    start: int,
    limit: int,
    ocr_mode: OCRMode,
) -> dict[str, Any]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    presentation = Presentation(io.BytesIO(content))
    total_slides = len(presentation.slides)
    end = min(total_slides, start + min(limit, 50) - 1)
    slides: list[dict[str, Any]] = []
    image_count = 0
    for slide_number in range(start, end + 1):
        slide = presentation.slides[slide_number - 1]
        texts: list[str] = []
        tables: list[dict[str, Any]] = []
        images: list[dict[str, Any]] = []
        for shape in _iter_presentation_shapes(slide.shapes):
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    texts.append(text[:_MAX_CELL_TEXT])
            if getattr(shape, "has_table", False):
                tables.append(
                    {
                        "rows": _bounded_rows(
                            (
                                (cell.text for cell in row.cells)
                                for row in shape.table.rows
                            ),
                        ),
                    },
                )
            if (
                shape.shape_type == MSO_SHAPE_TYPE.PICTURE
                and ocr_mode != "never"
                and image_count < _MAX_OCR_IMAGES_PER_CALL
            ):
                image_count += 1
                try:
                    images.append(
                        _ocr_image(
                            shape.image.blob,
                            source=(
                                f"PPTX 第{slide_number}页图片{image_count}"
                            ),
                        ),
                    )
                except Exception as exc:
                    images.append(
                        {
                            "source": (
                                f"PPTX 第{slide_number}页图片{image_count}"
                            ),
                            "error": str(exc),
                        },
                    )
        slides.append(
            {
                "slide": slide_number,
                "texts": texts,
                "tables": tables,
                "images": images,
            },
        )
    return {
        "format": "pptx",
        "start_slide": start,
        "end_slide": end,
        "total_slides": total_slides,
        "slides": slides,
        "next_start": end + 1 if end < total_slides else None,
    }


def _render_pdf_page(content: bytes, page_index: int) -> bytes:
    import pypdfium2 as pdfium

    document = pdfium.PdfDocument(content)
    try:
        page = document[page_index]
        try:
            image = page.render(scale=2.0).to_pil()
            output = io.BytesIO()
            image.save(output, format="PNG")
            return output.getvalue()
        finally:
            page.close()
    finally:
        document.close()


def _parse_pdf(
    content: bytes,
    start: int,
    limit: int,
    ocr_mode: OCRMode,
) -> dict[str, Any]:
    import pdfplumber

    with pdfplumber.open(io.BytesIO(content)) as document:
        total_pages = len(document.pages)
        end = min(total_pages, start + min(limit, 20) - 1)
        pages: list[dict[str, Any]] = []
        for page_number in range(start, end + 1):
            page = document.pages[page_number - 1]
            text = (page.extract_text() or "")[:50_000]
            raw_tables = page.extract_tables()[:_MAX_TABLES_PER_PAGE]
            tables = [
                {
                    "table": table_index,
                    "rows": _bounded_rows(table),
                }
                for table_index, table in enumerate(raw_tables, start=1)
            ]
            should_ocr = (
                ocr_mode == "always"
                or (ocr_mode == "auto" and len(text.strip()) < 20)
            )
            ocr: dict[str, Any] | None = None
            if should_ocr:
                try:
                    rendered = _render_pdf_page(
                        content,
                        page_number - 1,
                    )
                    ocr = _ocr_image(
                        rendered,
                        source=f"PDF 第{page_number}页",
                    )
                except Exception as exc:
                    ocr = {
                        "source": f"PDF 第{page_number}页",
                        "error": str(exc),
                    }
            pages.append(
                {
                    "page": page_number,
                    "text": text,
                    "tables": tables,
                    "ocr": ocr,
                },
            )
    return {
        "format": "pdf",
        "start_page": start,
        "end_page": end,
        "total_pages": total_pages,
        "pages": pages,
        "next_start": end + 1 if end < total_pages else None,
    }


def parse_attachment_content(
    content: bytes,
    suffix: str,
    sheet_name: str | None,
    start: int,
    limit: int,
    ocr_mode: str = "auto",
) -> dict[str, Any]:
    """Parse a bounded portion of one supported attachment."""
    normalized_suffix = suffix.strip().lower()
    normalized_ocr_mode = _normalize_ocr_mode(ocr_mode)
    normalized_start = max(1, int(start))
    normalized_limit = min(max(1, int(limit)), 500)

    if normalized_suffix == ".xlsx":
        return _parse_xlsx(
            content,
            sheet_name,
            normalized_start,
            normalized_limit,
            normalized_ocr_mode,
        )
    if normalized_suffix == ".xls":
        return _parse_xls(
            content,
            sheet_name,
            normalized_start,
            normalized_limit,
        )
    if normalized_suffix == ".csv":
        decoded, encoding = _decode_text_file(content)
        sample = decoded[:4096]
        try:
            dialect = csv.Sniffer().sniff(sample)
        except csv.Error:
            dialect = csv.excel
        rows = list(csv.reader(io.StringIO(decoded), dialect=dialect))
        selected = rows[
            normalized_start - 1:
            normalized_start - 1 + normalized_limit
        ]
        end = normalized_start + len(selected) - 1
        return {
            "format": "csv",
            "encoding": encoding,
            "delimiter": dialect.delimiter,
            "start_row": normalized_start,
            "end_row": end,
            "total_rows": len(rows),
            "rows": selected,
            "next_start": end + 1 if end < len(rows) else None,
        }
    if normalized_suffix in {".txt", ".md"}:
        decoded, encoding = _decode_text_file(content)
        lines = decoded.splitlines()
        selected = lines[
            normalized_start - 1:
            normalized_start - 1 + normalized_limit
        ]
        end = normalized_start + len(selected) - 1
        return {
            "format": normalized_suffix.removeprefix("."),
            "encoding": encoding,
            "start_line": normalized_start,
            "end_line": end,
            "total_lines": len(lines),
            "lines": selected,
            "next_start": end + 1 if end < len(lines) else None,
        }
    if normalized_suffix == ".docx":
        return _parse_docx(
            content,
            normalized_start,
            normalized_limit,
            normalized_ocr_mode,
        )
    if normalized_suffix == ".pptx":
        return _parse_pptx(
            content,
            normalized_start,
            normalized_limit,
            normalized_ocr_mode,
        )
    if normalized_suffix == ".pdf":
        return _parse_pdf(
            content,
            normalized_start,
            normalized_limit,
            normalized_ocr_mode,
        )
    if normalized_suffix in _IMAGE_SUFFIXES:
        if normalized_ocr_mode == "never":
            from PIL import Image, ImageOps

            with Image.open(io.BytesIO(content)) as original:
                image = ImageOps.exif_transpose(original)
                width, height = image.size
            return {
                "format": normalized_suffix.removeprefix("."),
                "image": {
                    "width": width,
                    "height": height,
                    "ocr": None,
                },
            }
        return {
            "format": normalized_suffix.removeprefix("."),
            "image": {
                "ocr": _ocr_image(content, source="上传图片"),
            },
        }
    raise ValueError(
        f"不支持的附件格式：{normalized_suffix or '未知'}",
    )
